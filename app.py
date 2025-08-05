import streamlit as st
from PIL import Image, ImageChops
import os
import io
import math
from pptx import Presentation
from pptx.util import Inches

PRELOADED_LOGO_DIR = "preloaded_logos"

def trim_whitespace(image):
    bg = Image.new(image.mode, image.size, (255, 255, 255, 0))
    diff = ImageChops.difference(image, bg)
    bbox = diff.getbbox()
    if bbox:
        return image.crop(bbox)
    return image

def resize_to_fill_5x2_box(image, cell_width_px, cell_height_px, buffer_ratio=0.7):
    box_ratio = 3 / 1
    max_box_width = int(cell_width_px * buffer_ratio)
    max_box_height = int(cell_height_px * buffer_ratio)

    # Dynamically adjust logo scale based on box size
    logo_scale = 0.92 if max_box_width < 250 else 0.85

    if max_box_width / box_ratio <= max_box_height:
        box_width = max_box_width
        box_height = int(max_box_width / box_ratio)
    else:
        box_height = max_box_height
        box_width = int(max_box_height * box_ratio)

    box_width = int(box_width * logo_scale)
    box_height = int(box_height * logo_scale)

    img_w, img_h = image.size
    img_ratio = img_w / img_h

    if img_ratio > (box_width / box_height):
        new_width = box_width
        new_height = int(box_width / img_ratio)
    else:
        new_height = box_height
        new_width = int(box_height * img_ratio)

    resized = image.resize((new_width, new_height), Image.LANCZOS)
    return resized, box_width, box_height

def create_logo_slide(prs, logos, canvas_width_in, canvas_height_in, logos_per_row):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    canvas_width_px = int(canvas_width_in * 96)
    canvas_height_px = int(canvas_height_in * 96)
    
    logo_count = len(logos)
    cols = logos_per_row if logos_per_row else max(1, round((logo_count / 1.5) ** 0.5 * (canvas_width_in / canvas_height_in) ** 0.3))
    rows = math.ceil(logo_count / cols)

    # Dynamic padding based on grid size
    base_padding = 0.85
    padding_boost = min(0.1, 0.2 / max(cols, rows))
    padding_ratio = base_padding + padding_boost

    cell_width = (canvas_width_px / cols) * padding_ratio
    cell_height = (canvas_height_px / rows) * padding_ratio

    left_margin = Inches((10 - canvas_width_in) / 2 + 0.1)
    top_margin = Inches((7.5 - canvas_height_in) / 2 + 0.1)

    for idx, logo in enumerate(logos):
        col = idx % cols
        row = idx // cols

        trimmed = trim_whitespace(logo)
        resized, box_w, box_h = resize_to_fill_5x2_box(trimmed, cell_width, cell_height)

        img_stream = io.BytesIO()
        resized.save(img_stream, format="PNG", dpi=(300, 300))
        img_stream.seek(0)

        x_offset = (cell_width - resized.width) / 2
        y_offset = (cell_height - resized.height) / 2

        left = left_margin + Inches((col * cell_width + x_offset) / 96)
        top = top_margin + Inches((row * cell_height + y_offset) / 96)

        slide.shapes.add_picture(
            img_stream,
            left,
            top,
            width=Inches(resized.width / 96),
            height=Inches(resized.height / 96)
        )

# --- Streamlit UI ---
st.title("Logo Grid PowerPoint Exporter")
st.markdown("Upload logos or use preloaded ones below:")

if not os.path.exists(PRELOADED_LOGO_DIR):
    os.makedirs(PRELOADED_LOGO_DIR)
preloaded_filenames = sorted([
    os.path.splitext(f)[0] for f in os.listdir(PRELOADED_LOGO_DIR)
    if f.lower().endswith((".png", ".jpg", ".jpeg", ".webp"))
], key=lambda x: x.lower())

uploaded_files = st.file_uploader("Upload logos", type=["png", "jpg", "jpeg", "webp"], accept_multiple_files=True)
selected_preloaded = st.multiselect("Select preloaded logos", options=preloaded_filenames)

canvas_width_in = st.number_input("Grid width (inches)", min_value=1.0, max_value=20.0, value=10.0)
canvas_height_in = st.number_input("Grid height (inches)", min_value=1.0, max_value=20.0, value=7.5)
logos_per_row = st.number_input("Logos per row (optional)", min_value=0, max_value=50, value=0)

if st.button("Generate PowerPoint"):
    logo_entries = []

    if uploaded_files:
        for f in uploaded_files:
            name = os.path.splitext(f.name)[0]
            image = Image.open(f).convert("RGBA")
            logo_entries.append((name.lower(), image))

    for name in selected_preloaded:
        for ext in [".png", ".jpg", ".jpeg", ".webp"]:
            path = os.path.join(PRELOADED_LOGO_DIR, name + ext)
            if os.path.exists(path):
                image = Image.open(path).convert("RGBA")
                logo_entries.append((name.lower(), image))
                break

    if not logo_entries:
        st.warning("Please upload or select logos.")
    else:
        logo_entries.sort(key=lambda x: x[0])
        images = [entry[1] for entry in logo_entries]
        prs = Presentation()
        create_logo_slide(prs, images, canvas_width_in, canvas_height_in,
                          logos_per_row if logos_per_row > 0 else None)
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        st.success("PowerPoint created!")
        st.download_button("Download .pptx", output, file_name="logo_grid.pptx")
